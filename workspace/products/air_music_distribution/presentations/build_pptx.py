from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
BLACK = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AIR_BLUE = RGBColor(0x00, 0x8C, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0x00, 0xC8, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, text, font_size=24, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False, bg_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def multiline_box(slide, left, top, width, height, lines, default_size=22,
                  default_color=WHITE, default_bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        text = line.get("text", "")
        size = line.get("size", default_size)
        color = line.get("color", default_color)
        bold = line.get("bold", default_bold)
        italic = line.get("italic", False)
        p.alignment = line.get("align", align)
        space_before = line.get("space_before", 0)
        if space_before:
            p.space_before = Pt(space_before)
        if not text:
            p.add_run()
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

multiline_box(s, Inches(0.5), Inches(1.0), Inches(9), Inches(4.5), [
    {"text": "You've worked hard", "size": 52, "bold": True, "color": WHITE},
    {"text": "to get here.", "size": 52, "bold": True, "color": WHITE, "space_before": 2},
    {"text": "Your music is out there.", "size": 36, "bold": False, "color": MID_GRAY, "space_before": 18},
    {"text": "But nothing is moving.", "size": 36, "bold": False, "color": MID_GRAY, "space_before": 2},
    {"text": "— Let's fix that.", "size": 28, "bold": True, "color": AIR_BLUE, "space_before": 24},
])

multiline_box(s, Inches(0.5), Inches(6.3), Inches(8), Inches(1), [
    {"text": "AIR Music Distribution  ·  air.io  ·  Trusted by Creators since 2010",
     "size": 14, "color": DARK_GRAY},
])

# Primavera label top right
box(s, Inches(9.8), Inches(0.3), Inches(3.2), Inches(0.6),
    "Primavera Sound 2026", 13, False, DARK_GRAY, PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 2 — THE FEELING
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

multiline_box(s, Inches(0.5), Inches(0.8), Inches(11), Inches(2.5), [
    {"text": '"I\'m not a beginner anymore.', "size": 44, "bold": True, "color": WHITE},
    {"text": 'But I don\'t feel like a professional either."',
     "size": 44, "bold": True, "color": WHITE, "space_before": 4},
])

multiline_box(s, Inches(0.5), Inches(3.8), Inches(10), Inches(2.5), [
    {"text": "You have releases. You have listeners.", "size": 24, "color": LIGHT_GRAY},
    {"text": "You're somewhere between indie and career.", "size": 24, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "And you're doing almost everything alone.", "size": 24, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "", "size": 8},
    {"text": "This is not about talent. It's about structure.",
     "size": 26, "bold": True, "color": AIR_BLUE, "space_before": 8},
])


# ─────────────────────────────────────────────
# SLIDE 3 — PAINS
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.4), Inches(10), Inches(0.7),
    "Sound familiar?", 34, True, WHITE)

pains = [
    ('"My release is stuck in processing — 3 days.\nNo one explains why."', Inches(0.5), Inches(1.4)),
    ('"A Content ID claim appeared overnight.\nSupport is completely silent."', Inches(6.9), Inches(1.4)),
    ('"My royalties don\'t add up.\nI can\'t get a straight answer."', Inches(0.5), Inches(3.8)),
    ('"I submitted for playlists.\nNothing. Zero feedback. Just silence."', Inches(6.9), Inches(3.8)),
]

for text, l, t in pains:
    rect(s, l - Inches(0.15), t - Inches(0.15), Inches(5.9), Inches(1.9), RGBColor(0x1A, 0x1A, 0x1A))
    box(s, l, t, Inches(5.6), Inches(1.7), text, 19, False, LIGHT_GRAY)

box(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.6),
    "You're not alone. And this is not how it should work.",
    16, True, AIR_BLUE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 4 — ROOT CAUSE
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

multiline_box(s, Inches(0.5), Inches(0.8), Inches(11), Inches(2), [
    {"text": "The distribution industry was built",
     "size": 40, "bold": True, "color": WHITE},
    {"text": "for volume — not for artists.", "size": 40, "bold": True, "color": AIR_BLUE, "space_before": 4},
])

multiline_box(s, Inches(0.5), Inches(3.2), Inches(10), Inches(3.2), [
    {"text": "Most platforms make money whether your career grows or not.",
     "size": 24, "color": LIGHT_GRAY},
    {"text": "They have no reason to pick up the phone.", "size": 24, "color": LIGHT_GRAY, "space_before": 8},
    {"text": "No reason to care about your specific release.", "size": 24, "color": LIGHT_GRAY, "space_before": 8},
    {"text": "No reason to explain anything.", "size": 24, "color": LIGHT_GRAY, "space_before": 8},
    {"text": "", "size": 6},
    {"text": "You're not a priority. You're a number.", "size": 28, "bold": True, "color": WHITE, "space_before": 16},
])


# ─────────────────────────────────────────────
# SLIDE 5 — INTRODUCING AIR
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

multiline_box(s, Inches(0.5), Inches(0.8), Inches(11), Inches(2), [
    {"text": "AIR is not a platform.", "size": 52, "bold": True, "color": WHITE},
    {"text": "AIR is a team.", "size": 52, "bold": True, "color": AIR_BLUE, "space_before": 8},
])

multiline_box(s, Inches(0.5), Inches(3.2), Inches(10), Inches(2.8), [
    {"text": "14+ years on the market  ·  Trusted by Creators since 2010",
     "size": 20, "color": MID_GRAY},
    {"text": "", "size": 6},
    {"text": "Official partner:", "size": 18, "color": MID_GRAY, "space_before": 12},
    {"text": "YouTube  ·  Google  ·  TikTok  ·  Microsoft",
     "size": 26, "bold": True, "color": WHITE, "space_before": 4},
    {"text": "", "size": 6},
    {"text": "Distribution to 100+ streaming platforms worldwide",
     "size": 20, "color": LIGHT_GRAY, "space_before": 12},
    {"text": "Spotify · Apple Music · YouTube Music · TikTok · Instagram Reels · and more",
     "size": 16, "color": MID_GRAY, "space_before": 4},
])


# ─────────────────────────────────────────────
# SLIDE 6 — WHAT YOU GET
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.4), Inches(11), Inches(0.7),
    "Your AIR distribution includes:", 30, True, WHITE)

items = [
    ("👤  Personal Manager", "A real human who knows your music and goals"),
    ("⚡  4-Hour Response", "Not days. Not a ticket. A person."),
    ("🎵  Playlist Pitching", "Editorial playlists. Included. No extra fee."),
    ("📱  UGC Monetization", "Earn every time your track is used on TikTok, Reels, Shorts"),
    ("💰  70% Revenue Share", "Zero upfront. Zero hidden fees. We earn when you earn."),
    ("🔐  100% Your Rights", "Always. No exceptions."),
]

cols = [(Inches(0.5), Inches(1.3)), (Inches(7.0), Inches(1.3)),
        (Inches(0.5), Inches(3.0)), (Inches(7.0), Inches(3.0)),
        (Inches(0.5), Inches(4.7)), (Inches(7.0), Inches(4.7))]

for (title, sub), (l, t) in zip(items, cols):
    rect(s, l - Inches(0.1), t - Inches(0.1), Inches(5.8), Inches(1.4),
         RGBColor(0x16, 0x16, 0x16))
    box(s, l, t, Inches(5.5), Inches(0.55), title, 19, True, WHITE)
    box(s, l, t + Inches(0.5), Inches(5.5), Inches(0.7), sub, 15, False, MID_GRAY)


# ─────────────────────────────────────────────
# SLIDE 7 — PERSONAL MANAGER
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

rect(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(2.1), RGBColor(0x0A, 0x1A, 0x2E))
multiline_box(s, Inches(0.7), Inches(0.9), Inches(11.8), Inches(1.8), [
    {"text": '"When something goes wrong at 2am —', "size": 32, "bold": True, "color": WHITE},
    {"text": 'you have someone to call."', "size": 32, "bold": True, "color": AIR_BLUE, "space_before": 4},
])

multiline_box(s, Inches(0.5), Inches(3.2), Inches(11.5), Inches(3.5), [
    {"text": "Not a template reply.", "size": 22, "color": LIGHT_GRAY},
    {"text": "Not a support ticket that disappears for a week.", "size": 22, "color": LIGHT_GRAY, "space_before": 6},
    {"text": "", "size": 6},
    {"text": "A real person who:", "size": 22, "bold": True, "color": WHITE, "space_before": 10},
    {"text": "  →  Knows your catalog and your upcoming release",
     "size": 20, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "  →  Checks your track before it gets flagged",
     "size": 20, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "  →  Explains exactly why royalties look the way they do",
     "size": 20, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "  →  Goes to bat for you with the platforms",
     "size": 20, "color": LIGHT_GRAY, "space_before": 4},
    {"text": "", "size": 6},
    {"text": "We check, explain, and take responsibility.",
     "size": 22, "bold": True, "color": AIR_BLUE, "space_before": 8},
])


# ─────────────────────────────────────────────
# SLIDE 8 — NUMBERS
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.35), Inches(11), Inches(0.65),
    "The numbers that matter", 30, True, WHITE)

stats = [
    ("14+", "Years on market"),
    ("100+", "Streaming platforms"),
    ("4 hrs", "Avg. response time"),
    ("70%", "Your revenue share"),
    ("0", "Upfront payments"),
    ("100%", "Rights stay with you"),
]

positions = [
    (Inches(0.5), Inches(1.3)),
    (Inches(4.7), Inches(1.3)),
    (Inches(8.9), Inches(1.3)),
    (Inches(0.5), Inches(4.0)),
    (Inches(4.7), Inches(4.0)),
    (Inches(8.9), Inches(4.0)),
]

for (num, label), (l, t) in zip(stats, positions):
    rect(s, l, t, Inches(3.6), Inches(2.2), RGBColor(0x10, 0x10, 0x10))
    box(s, l + Inches(0.15), t + Inches(0.15), Inches(3.3), Inches(1.1),
        num, 52, True, AIR_BLUE, PP_ALIGN.CENTER)
    box(s, l + Inches(0.15), t + Inches(1.3), Inches(3.3), Inches(0.7),
        label, 16, False, MID_GRAY, PP_ALIGN.CENTER)

box(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
    "2-year contract  ·  Flexible exit  ·  No penalties",
    15, False, DARK_GRAY, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 9 — HOW IT WORKS
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.35), Inches(11), Inches(0.65),
    "From conversation to career — 3 steps", 30, True, WHITE)

steps = [
    ("01", "Apply", "Talk to a real person about your catalog and your goals.\nNo forms that vanish. A conversation."),
    ("02", "Get Your Manager", "Your dedicated manager reviews your tracks,\nsets up distribution, handles playlist pitching."),
    ("03", "Grow — Not Alone", "Monthly reports. Transparent analytics.\nSomeone in your corner, every step."),
]

for i, (num, title, desc) in enumerate(steps):
    l = Inches(0.5 + i * 4.3)
    t = Inches(1.4)
    rect(s, l, t, Inches(3.9), Inches(5.0), RGBColor(0x10, 0x10, 0x10))
    box(s, l + Inches(0.2), t + Inches(0.2), Inches(3.5), Inches(1.0),
        num, 48, True, AIR_BLUE)
    box(s, l + Inches(0.2), t + Inches(1.4), Inches(3.5), Inches(0.6),
        title, 22, True, WHITE)
    box(s, l + Inches(0.2), t + Inches(2.2), Inches(3.5), Inches(2.5),
        desc, 16, False, MID_GRAY)

    if i < 2:
        box(s, l + Inches(4.1), Inches(3.8), Inches(0.3), Inches(0.5),
            "→", 28, True, AIR_BLUE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 10 — HONEST COMPARISON
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.35), Inches(11), Inches(0.65),
    "What most distributors offer vs. what AIR offers", 28, True, WHITE)

# Table header
rect(s, Inches(0.5), Inches(1.25), Inches(4.3), Inches(0.55), RGBColor(0x22, 0x22, 0x22))
rect(s, Inches(4.85), Inches(1.25), Inches(3.7), Inches(0.55), RGBColor(0x22, 0x22, 0x22))
rect(s, Inches(8.6), Inches(1.25), Inches(4.2), Inches(0.55), RGBColor(0x00, 0x55, 0xAA))

box(s, Inches(0.6), Inches(1.3), Inches(4.0), Inches(0.45), "Category", 16, True, MID_GRAY)
box(s, Inches(4.95), Inches(1.3), Inches(3.5), Inches(0.45), "Most Distributors", 16, True, MID_GRAY)
box(s, Inches(8.7), Inches(1.3), Inches(4.0), Inches(0.45), "AIR", 16, True, WHITE)

rows = [
    ("Support", "Email · 5–7 days", "Personal manager · 4 hours"),
    ("Playlist Pitching", "Paid add-on", "Included"),
    ("Business Model", "Fixed fee — you pay\nwhether you earn or not", "Revenue share —\nwe earn when you earn"),
    ("Your Rights", "Varies", "100% yours, always"),
    ("Exit Terms", "Often locked in", "Flexible, no penalties"),
]

for i, (cat, bad, good) in enumerate(rows):
    t = Inches(1.9 + i * 0.95)
    bg_row = RGBColor(0x14, 0x14, 0x14) if i % 2 == 0 else BLACK
    rect(s, Inches(0.5), t, Inches(4.3), Inches(0.85), bg_row)
    rect(s, Inches(4.85), t, Inches(3.7), Inches(0.85), bg_row)
    rect(s, Inches(8.6), t, Inches(4.2), Inches(0.85), RGBColor(0x05, 0x18, 0x30))
    box(s, Inches(0.65), t + Inches(0.05), Inches(3.9), Inches(0.75), cat, 16, True, WHITE)
    box(s, Inches(4.95), t + Inches(0.05), Inches(3.4), Inches(0.75), bad, 14, False, MID_GRAY)
    box(s, Inches(8.7), t + Inches(0.05), Inches(4.0), Inches(0.75), good, 14, True, ACCENT)

box(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.5),
    "If you don't grow — we don't grow. That's the model.",
    15, True, AIR_BLUE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 11 — WHO THIS IS FOR
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

box(s, Inches(0.5), Inches(0.35), Inches(11), Inches(0.65),
    "AIR is built for artists like you", 30, True, WHITE)

for_items = [
    "You have an audience — 1K to 500K monthly listeners",
    "You have catalog — at least 5 releases",
    "You're tired of fighting your distributor instead of making music",
    "You want structure and a real partner — not just a platform",
]

for i, item in enumerate(for_items):
    t = Inches(1.25 + i * 0.82)
    rect(s, Inches(0.5), t, Inches(0.5), Inches(0.62), RGBColor(0x00, 0x77, 0xCC))
    box(s, Inches(0.6), t + Inches(0.05), Inches(0.3), Inches(0.5), "✓", 22, True, WHITE, PP_ALIGN.CENTER)
    box(s, Inches(1.15), t + Inches(0.07), Inches(11), Inches(0.5), item, 20, False, LIGHT_GRAY)

rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.04), RGBColor(0x33, 0x33, 0x33))

box(s, Inches(0.5), Inches(4.9), Inches(11), Inches(0.5),
    "AIR is NOT for you if:", 20, True, MID_GRAY)

not_for = [
    "You're uploading your first track",
    "You want the cheapest option with zero support",
    "You're fine being ticket #4821 in a queue",
]

for i, item in enumerate(not_for):
    t = Inches(5.5 + i * 0.55)
    box(s, Inches(0.5), t, Inches(12), Inches(0.5),
        f"  ✗  {item}", 17, False, DARK_GRAY)

box(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
    "We're selective because we take every artist seriously.",
    15, True, AIR_BLUE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 12 — PRIMAVERA OFFER
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

rect(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.05), RGBColor(0x05, 0x18, 0x30))
multiline_box(s, Inches(0.7), Inches(0.4), Inches(11.8), Inches(0.85), [
    {"text": "You're at Primavera. That means you're serious about your career.",
     "size": 24, "bold": True, "color": WHITE},
])

box(s, Inches(0.5), Inches(1.6), Inches(11), Inches(0.55),
    "Here's what we're offering — right now, right here:", 22, True, LIGHT_GRAY)

offers = [
    ("🎯", "Free Track Safety Check",
     "We audit your current distribution setup — Content ID risks,\nroyalty leaks, metadata errors. Free. No strings."),
    ("💬", "Personal Consultation — Today",
     "Talk to our team. Not a form, not an email.\nA real conversation — right here at Primavera."),
    ("⚡", "Fast-Track Onboarding",
     "Artists who connect with us at Primavera\nskip the standard queue. Priority onboarding."),
]

for i, (icon, title, desc) in enumerate(offers):
    l = Inches(0.5 + i * 4.3)
    t = Inches(2.4)
    rect(s, l, t, Inches(3.9), Inches(3.8), RGBColor(0x0D, 0x1F, 0x35))
    box(s, l + Inches(0.2), t + Inches(0.2), Inches(3.5), Inches(0.75),
        icon, 36, False, WHITE, PP_ALIGN.CENTER)
    box(s, l + Inches(0.2), t + Inches(1.1), Inches(3.5), Inches(0.65),
        title, 18, True, ACCENT)
    box(s, l + Inches(0.2), t + Inches(1.85), Inches(3.5), Inches(1.8),
        desc, 14, False, MID_GRAY)


# ─────────────────────────────────────────────
# SLIDE 13 — CTA / CLOSING
# ─────────────────────────────────────────────
s = add_slide()
bg(s, BLACK)
rect(s, Inches(0), Inches(0), Inches(0.12), H, AIR_BLUE)

rect(s, Inches(0), Inches(2.8), W, Inches(0.06), AIR_BLUE)

multiline_box(s, Inches(0.5), Inches(0.6), Inches(12), Inches(2.0), [
    {"text": "Ready to stop fighting", "size": 52, "bold": True, "color": WHITE},
    {"text": "your distributor?", "size": 52, "bold": True, "color": AIR_BLUE, "space_before": 4},
])

multiline_box(s, Inches(0.5), Inches(3.3), Inches(11.5), Inches(2.8), [
    {"text": "→  Scan the QR code below — book your Free Track Safety Check",
     "size": 22, "color": LIGHT_GRAY},
    {"text": "→  Come talk to us — our team is here today",
     "size": 22, "color": LIGHT_GRAY, "space_before": 10},
    {"text": "→  air.io — see everything for yourself",
     "size": 22, "color": LIGHT_GRAY, "space_before": 10},
])

multiline_box(s, Inches(0.5), Inches(6.3), Inches(11), Inches(0.8), [
    {"text": "AIR Music Distribution  ·  air.io  ·  \"We earn when you earn.\"",
     "size": 16, "color": DARK_GRAY},
])

# QR placeholder box
rect(s, Inches(10.8), Inches(4.8), Inches(2.1), Inches(2.1), RGBColor(0x18, 0x18, 0x18))
box(s, Inches(10.8), Inches(5.55), Inches(2.1), Inches(0.5),
    "QR Code\nair.io", 12, False, MID_GRAY, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out = "/home/user/Nataliia_Borska/workspace/products/air_music_distribution/presentations/AIR_Primavera2026.pptx"
prs.save(out)
print(f"Saved: {out}")
